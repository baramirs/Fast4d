"""PDF / DOCX / PPTX writers — compose pages from export assets (no PIL crops)."""
from __future__ import annotations

from datetime import datetime
from pathlib import Path


def _cover_lines(rows: list[dict], title: str) -> list[str]:
    scans = sorted({r["scan"] for r in rows})
    n_calib = sum(1 for r in rows if r.get("kind") == "calib")
    n_maps = sum(1 for r in rows if r.get("kind") in ("strain", "stress"))
    n_rep = sum(1 for r in rows if r.get("kind") == "report")
    return [
        title,
        datetime.now().strftime("%Y-%m-%d %H:%M"),
        f"Scans: {', '.join(scans) if scans else '(none)'}",
        f"Panels: {len(rows)}  (calib={n_calib}, maps={n_maps}, reports={n_rep})",
        "Each panel is a full figure — not cropped from a collage.",
    ]


def write_pdf(rows: list[dict], out_path: Path, *, title: str = "Fast4D Report") -> Path:
    """One panel per page (portrait)."""
    import matplotlib
    matplotlib.use("Agg", force=False)
    import matplotlib.image as mpimg
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    lines = _cover_lines(rows, title)
    with PdfPages(out_path) as pdf:
        fig = plt.figure(figsize=(8.5, 11))
        y = 0.72
        for i, line in enumerate(lines):
            size = 20 if i == 0 else 11
            weight = "bold" if i == 0 else "normal"
            fig.text(0.5, y, line, ha="center", va="center",
                     fontsize=size, fontweight=weight)
            y -= 0.06
        pdf.savefig(fig)
        plt.close(fig)

        current_scan = None
        for row in rows:
            if row["scan"] != current_scan:
                current_scan = row["scan"]
                fig = plt.figure(figsize=(8.5, 11))
                fig.text(0.5, 0.55, current_scan, ha="center", fontsize=18,
                         fontweight="bold")
                fig.text(0.5, 0.48, "Scan section", ha="center", fontsize=12)
                pdf.savefig(fig)
                plt.close(fig)
            fig, ax = plt.subplots(figsize=(8.5, 11))
            img = mpimg.imread(str(row["path"]))
            ax.imshow(img)
            ax.set_axis_off()
            fig.suptitle(
                f"{row.get('section', '')}: {row['title']}", fontsize=11, y=0.98)
            fig.tight_layout(rect=(0, 0, 1, 0.96))
            pdf.savefig(fig)
            plt.close(fig)
    return out_path


def write_docx(rows: list[dict], out_path: Path, *, title: str = "Fast4D Report") -> Path:
    try:
        from docx import Document
        from docx.shared import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(
            "python-docx is required for DOCX export. "
            "Install with: pip install python-docx"
        ) from exc

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    for i, line in enumerate(_cover_lines(rows, title)):
        if i == 0:
            doc.add_heading(line, level=0)
        else:
            doc.add_paragraph(line)

    current_scan = None
    current_section = None
    for row in rows:
        if row["scan"] != current_scan:
            current_scan = row["scan"]
            current_section = None
            doc.add_heading(current_scan, level=1)
        sec = row.get("section") or row.get("label_title") or ""
        if sec != current_section:
            current_section = sec
            doc.add_heading(sec, level=2)
        doc.add_heading(str(row.get("channel", "")), level=3)
        doc.add_picture(str(row["path"]), width=Inches(6.0))
        p = doc.add_paragraph(row["title"])
        for run in p.runs:
            run.font.size = Pt(9)
    doc.save(out_path)
    return out_path


def write_pptx(rows: list[dict], out_path: Path, *, title: str = "Fast4D Report",
               template: Path | None = None) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(template)) if template else Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    def _title_slide(text: str, subtitle: str = "") -> None:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.2))
        tf = box.text_frame
        tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.bold = True
        run.font.size = Pt(28)
        if subtitle:
            box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12), Inches(1))
            tf2 = box2.text_frame
            tf2.clear()
            run2 = tf2.paragraphs[0].add_run()
            run2.text = subtitle
            run2.font.size = Pt(14)

    cover = _cover_lines(rows, title)
    _title_slide(cover[0], "\n".join(cover[1:]))

    current_scan = None
    for row in rows:
        if row["scan"] != current_scan:
            current_scan = row["scan"]
            _title_slide(current_scan, "Scan section")
        slide = prs.slides.add_slide(blank)
        tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.5), Inches(0.5))
        tf = tbox.text_frame
        tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = f"[{row.get('section', '')}] {row['title']}"
        run.font.bold = True
        run.font.size = Pt(16)
        # Full panel — never a vertical strip crop of a collage
        slide.shapes.add_picture(
            str(row["path"]), Inches(0.6), Inches(0.7), height=Inches(6.4))

    prs.save(str(out_path))
    return out_path
