"""Build a PowerPoint summary from saved Fast4D calibration images.

Usage examples:

    python tools/export_calibration_pptx.py ^
      --summary "D:\\results\\summary" ^
      --out "D:\\results\\calibrations_report.pptx"

    python tools/export_calibration_pptx.py ^
      --summary "D:\\results\\summary" ^
      --template "D:\\templates\\ASU_template.pptx" ^
      --out "D:\\results\\calibrations_report.pptx"

Expected input layout (created by Fast4D Save):

    summary/
      calibrations/
        origin/*_origin.png
        ellipse/*_ellipse.png
        q_pixel/*_q_pixel.png
        basis/*_basis.png

The script creates one slide per calibration. For basis, it also splits each
composite basis PNG into 3 equal panels and creates additional slides for those
panels. This is intentionally independent from the Fast4D GUI.
"""
from __future__ import annotations

import argparse
import math
import re
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


STEP_LABELS = {
    "origin": "Origin",
    "ellipse": "Ellipse",
    "q_pixel": "Q Pixel Size",
    "basis": "Basis",
}

MAP_SPECS = [
    ("strain_without_roi", "Strain without ROI", ["exx", "eyy", "exy", "orientation"]),
    ("strain_with_roi", "Strain ROI", ["exx", "eyy", "exy", "orientation"]),
    ("stress_without_roi", "Stress without ROI", ["sxx", "syy", "sxy"]),
    ("stress_with_roi", "Stress ROI", ["sxx", "syy", "sxy"]),
]


def _safe_name(text: str) -> str:
    return re.sub(r"[^A-Za-z0-9._ -]+", "_", str(text)).strip() or "image"


def _blank_layout(prs: Presentation):
    # Built-in blank is usually index 6. Fall back to the last layout.
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]


def _add_title(slide, title: str, *, left=0.25, top=0.18, width=2.0, height=0.55):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(24)
    return box


def _image_label(path: Path, step: str) -> str:
    stem = path.stem
    suffix = f"_{step}"
    if stem.endswith(suffix):
        stem = stem[: -len(suffix)]
    stem = stem.replace("_basis_panel1", "").replace("_basis_panel2", "").replace("_basis_panel3", "")
    return stem


def _add_image_grid(slide, images: list[Path], *, step: str, title: str) -> None:
    _add_title(slide, title)
    if not images:
        box = slide.shapes.add_textbox(Inches(2.6), Inches(2.8), Inches(7.0), Inches(0.4))
        box.text_frame.text = "No images found."
        return

    # 16:9 default page: 13.333 x 7.5 in. Leave a title/logo margin on the left.
    left0, top0 = 2.55, 0.25
    area_w, area_h = 10.45, 6.85
    n = len(images)
    cols = 3 if n > 2 else max(1, n)
    rows = max(1, math.ceil(n / cols))
    gap_x, gap_y = 0.25, 0.28
    cell_w = (area_w - gap_x * (cols - 1)) / cols
    cell_h = (area_h - gap_y * (rows - 1)) / rows

    for i, img in enumerate(images):
        r, c = divmod(i, cols)
        x = left0 + c * (cell_w + gap_x)
        y = top0 + r * (cell_h + gap_y)
        label_h = 0.28
        label = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(cell_w), Inches(label_h))
        tf = label.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = 1  # center
        run = p.add_run()
        run.text = _image_label(img, step)
        run.font.bold = True
        run.font.size = Pt(14)

        pic_top = y + label_h
        pic_h = max(0.5, cell_h - label_h)
        # python-pptx preserves aspect ratio if only width is provided. If too tall,
        # retry by height.
        shp = slide.shapes.add_picture(str(img), Inches(x), Inches(pic_top), width=Inches(cell_w))
        if shp.height > Inches(pic_h):
            slide.shapes._spTree.remove(shp._element)
            shp = slide.shapes.add_picture(str(img), Inches(x), Inches(pic_top), height=Inches(pic_h))
        # Center within cell.
        shp.left = Inches(x) + max(0, int((Inches(cell_w) - shp.width) / 2))


def _add_wide_image(slide, image: Path, *, title: str) -> None:
    _add_title(slide, title, width=5.0)
    if not image.is_file():
        box = slide.shapes.add_textbox(Inches(2.6), Inches(2.8), Inches(7.0), Inches(0.4))
        box.text_frame.text = f"Missing image: {image}"
        return
    x, y, w, h = 1.0, 0.9, 11.3, 6.25
    shp = slide.shapes.add_picture(str(image), Inches(x), Inches(y), width=Inches(w))
    if shp.height > Inches(h):
        slide.shapes._spTree.remove(shp._element)
        shp = slide.shapes.add_picture(str(image), Inches(x), Inches(y), height=Inches(h))
    shp.left = Inches(x) + max(0, int((Inches(w) - shp.width) / 2))


def _find_step_images(calib_dir: Path, step: str) -> list[Path]:
    d = calib_dir / step
    if not d.is_dir():
        return []
    # Keep per-file calibration images, not metric plots like ellipse_a_px.png.
    return sorted(p for p in d.glob(f"*_{step}.png") if p.is_file())


def _metric_plot_images(calib_dir: Path, step: str) -> list[Path]:
    """Calibration value plots already exported by Fast4D (exclude per-file images)."""
    d = calib_dir / step
    if not d.is_dir():
        return []
    return sorted(
        p for p in d.glob("*.png")
        if p.is_file() and not p.name.endswith(f"_{step}.png")
    )


def split_basis_images(calib_dir: Path, *, panels: int = 3) -> dict[int, list[Path]]:
    """Split each basis composite PNG into N equal vertical panels."""
    basis_dir = calib_dir / "basis"
    out_dir = basis_dir / "split"
    split: dict[int, list[Path]] = {i: [] for i in range(1, panels + 1)}
    if not basis_dir.is_dir():
        return split
    out_dir.mkdir(parents=True, exist_ok=True)
    for img_path in sorted(basis_dir.glob("*_basis.png")):
        try:
            with Image.open(img_path) as im:
                w, h = im.size
                for i in range(panels):
                    x0 = int(round(w * i / panels))
                    x1 = int(round(w * (i + 1) / panels))
                    crop = im.crop((x0, 0, x1, h))
                    out = out_dir / f"{img_path.stem}_panel{i + 1}.png"
                    crop.save(out)
                    split[i + 1].append(out)
        except Exception as exc:
            print(f"[warn] basis split skipped for {img_path}: {exc}")
    return split


def _iter_scan_figure_dirs(summary: Path) -> list[Path]:
    """Find <batch>/<scan>/figures folders from the summary folder."""
    root = summary.parent
    return sorted(p for p in root.glob("*/figures") if p.is_dir())


def split_map_figures(summary: Path) -> dict[tuple[str, str], list[Path]]:
    """Split saved composite strain/stress figures into per-map panel images.

    Fast4D saves per-scan composite figures such as
    ``<scan>/figures/strain_with_roi.png``. This routine crops those composites
    into equal panels so a PPT slide can be made for each channel.
    """
    out: dict[tuple[str, str], list[Path]] = {}
    split_root = summary / "maps_split"
    split_root.mkdir(parents=True, exist_ok=True)
    for fig_dir in _iter_scan_figure_dirs(summary):
        scan = fig_dir.parent.name
        for fig_key, _label, channels in MAP_SPECS:
            src = fig_dir / f"{fig_key}.png"
            if not src.is_file():
                continue
            # Fast4D strain composites are py4DSTEM show_strain(layout="horizontal"):
            # one row of 5 panels → εxx, εyy, εxy, θ, legend(g1/g2). The 4 data panels
            # occupy the first ~80% of the width (the trailing legend panel is excluded).
            # Stress composites are 1x3 (build_stress_maps_figure) → equal vertical strips.
            try:
                with Image.open(src) as im:
                    w, h = im.size
                    if len(channels) == 4:
                        data_w = round(w * 0.80)        # drop the g1/g2 legend panel
                        boxes = [
                            (round(data_w * i / 4), 0, round(data_w * (i + 1) / 4), h)
                            for i in range(4)
                        ]
                    else:
                        boxes = [
                            (int(round(w * i / len(channels))), 0,
                             int(round(w * (i + 1) / len(channels))), h)
                            for i in range(len(channels))
                        ]
                    for ch, box in zip(channels, boxes):
                        d = split_root / fig_key / ch
                        d.mkdir(parents=True, exist_ok=True)
                        out_path = d / f"{_safe_name(scan)}_{fig_key}_{ch}.png"
                        im.crop(box).save(out_path)
                        out.setdefault((fig_key, ch), []).append(out_path)
            except Exception as exc:
                print(f"[warn] map split skipped for {src}: {exc}")
    return out


def create_calibration_trend_plots(calib_dir: Path) -> list[Path]:
    """Create clean file-vs-value trend plots from calibration_values_all.csv."""
    csv = calib_dir / "calibration_values_all.csv"
    out_dir = calib_dir / "trends"
    if not csv.is_file():
        return []
    out_dir.mkdir(parents=True, exist_ok=True)
    df = pd.read_csv(csv)
    if df.empty or "scan" not in df:
        return []
    scans = [str(s) for s in df["scan"]]
    x = list(range(len(scans)))
    written: list[Path] = []

    def save_plot(filename: str, title: str, series: list[tuple[str, str, str]]) -> None:
        fig, ax = plt.subplots(figsize=(10.8, 4.6), constrained_layout=True)
        any_series = False
        for col, label, color in series:
            if col not in df:
                continue
            y = pd.to_numeric(df[col], errors="coerce")
            if not y.notna().any():
                continue
            ax.plot(x, y, "o-", label=label, color=color, alpha=0.62, lw=1.8, ms=5)
            any_series = True
        if not any_series:
            plt.close(fig)
            return
        ax.set_xticks(x)
        ax.set_xticklabels(scans, rotation=55, ha="right", fontsize=8)
        ax.set_xlabel("files")
        ax.set_ylabel("data")
        ax.set_title(title)
        ax.grid(alpha=0.3)
        ax.legend()
        out = out_dir / filename
        fig.savefig(out, dpi=160, bbox_inches="tight")
        plt.close(fig)
        written.append(out)

    save_plot(
        "q_pixel_guess_vs_fitted.png",
        "Q-pixel guess vs fitted",
        [
            ("q_pixel_guess_Ainv_per_px", "guess", "#1565C0"),
            ("q_pixel_fitted_Ainv_per_px", "fitted", "#C62828"),
        ],
    )
    save_plot(
        "ellipse_a_b_values.png",
        "Ellipse fitted values (a, b)",
        [
            ("ellipse_a_px", "a fitted", "#2E7D32"),
            ("ellipse_b_px", "b fitted", "#EF6C00"),
        ],
    )
    save_plot(
        "ellipse_guess_vs_fitted.png",
        "Ellipse guess vs fitted (when guess columns exist)",
        [
            ("ellipse_a_guess_px", "a guess", "#1565C0"),
            ("ellipse_a_px", "a fitted", "#C62828"),
            ("ellipse_b_guess_px", "b guess", "#00838F"),
            ("ellipse_b_px", "b fitted", "#EF6C00"),
        ],
    )
    save_plot(
        "ellipse_ratio_theta.png",
        "Ellipse ratio / theta",
        [
            ("ellipse_a_over_b", "a/b", "#6A1B9A"),
            ("ellipse_theta_deg", "theta deg", "#00838F"),
        ],
    )
    return written


def build_pptx(summary: Path, out: Path, *, template: Path | None = None,
               split_basis: bool = True, include_maps: bool = True,
               include_trends: bool = True) -> Path:
    calib_dir = summary / "calibrations"
    if not calib_dir.is_dir():
        raise FileNotFoundError(f"Missing calibrations folder: {calib_dir}")

    prs = Presentation(str(template)) if template else Presentation()
    layout = _blank_layout(prs)

    for step in ("origin", "ellipse", "q_pixel", "basis"):
        images = _find_step_images(calib_dir, step)
        slide = prs.slides.add_slide(layout)
        _add_image_grid(slide, images, step=step, title=STEP_LABELS.get(step, step))

    if include_trends:
        for plot in create_calibration_trend_plots(calib_dir):
            slide = prs.slides.add_slide(layout)
            _add_wide_image(slide, plot, title=f"Calibration Trend - {plot.stem}")

        # Also include the individual metric PNGs Fast4D exported per calibration.
        for step in ("origin", "ellipse", "q_pixel", "basis"):
            plots = _metric_plot_images(calib_dir, step)
            if plots:
                slide = prs.slides.add_slide(layout)
                _add_image_grid(slide, plots, step=step, title=f"{STEP_LABELS.get(step, step)} Values")

    if split_basis:
        split = split_basis_images(calib_dir, panels=3)
        for panel_idx, images in split.items():
            if not images:
                continue
            slide = prs.slides.add_slide(layout)
            _add_image_grid(
                slide, images, step="basis",
                title=f"Basis - Panel {panel_idx}",
            )

    if include_maps:
        maps = split_map_figures(summary)
        for fig_key, title, channels in MAP_SPECS:
            for ch in channels:
                images = maps.get((fig_key, ch), [])
                if not images:
                    continue
                slide = prs.slides.add_slide(layout)
                _add_image_grid(slide, images, step=ch, title=f"{title} - {ch}")

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description="Build PPTX from Fast4D calibration summary images.")
    ap.add_argument("--summary", required=True, help="Path to summary folder.")
    ap.add_argument("--out", required=True, help="Output .pptx path.")
    ap.add_argument("--template", default="", help="Optional .pptx template.")
    ap.add_argument("--no-split-basis", action="store_true",
                    help="Do not create extra slides from split basis panels.")
    ap.add_argument("--no-maps", action="store_true",
                    help="Do not add strain/stress map slides.")
    ap.add_argument("--no-trends", action="store_true",
                    help="Do not add calibration trend/value plot slides.")
    ns = ap.parse_args()

    summary = Path(ns.summary).expanduser().resolve()
    out = Path(ns.out).expanduser().resolve()
    template = Path(ns.template).expanduser().resolve() if ns.template else None
    if template and not template.is_file():
        raise FileNotFoundError(f"Template not found: {template}")

    written = build_pptx(summary, out, template=template,
                         split_basis=not ns.no_split_basis,
                         include_maps=not ns.no_maps,
                         include_trends=not ns.no_trends)
    print(f"Wrote {written}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
